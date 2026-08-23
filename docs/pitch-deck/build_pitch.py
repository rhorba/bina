# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- Design tokens (from Bina ui-designer skill) ----------
NAVY      = RGBColor(0x1E, 0x3A, 0x5F)   # primary steel blue
NAVY_DK   = RGBColor(0x14, 0x27, 0x40)   # darker navy for title bg
NAVY_MID  = RGBColor(0x33, 0x50, 0x76)
ORANGE    = RGBColor(0xE0, 0x7B, 0x39)   # construction orange
CONCRETE  = RGBColor(0xF4, 0xF2, 0xEF)   # bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1B, 0x24, 0x30)   # foreground text
MUTED     = RGBColor(0x6B, 0x72, 0x80)
BORDER    = RGBColor(0xDD, 0xD8, 0xD0)
RED       = RGBColor(0xC0, 0x39, 0x2B)
AMBER     = RGBColor(0xC9, 0x8A, 0x1E)
GREEN     = RGBColor(0x3F, 0x8F, 0x5F)
SURFACE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "IBM Plex Sans"
FONT_AR = "Noto Kufi Arabic"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

def add_slide(bg=CONCRETE):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s

def rect(s, x, y, w, h, color, line=False, line_color=None, line_w=Pt(1), shadow_off=True):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or color
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    if shadow_off:
        shp.shadow.inherit = False
    return shp

def rrect(s, x, y, w, h, color, radius=0.08, line=False, line_color=None, line_w=Pt(1)):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or color
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def textbox(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
            font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True, italic=False,
            spacing_after=0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if spacing_after:
            p.space_after = Pt(spacing_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb

def multi_run_para(tf, runs, align=PP_ALIGN.LEFT, line_spacing=1.0, first=False, space_after=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_after:
        p.space_after = Pt(space_after)
    for text, size, color, bold, font in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return p

def bullets(s, x, y, w, h, items, size=15, color=INK, bullet_color=ORANGE, gap=10, font=FONT,
            bold_lead=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        r1 = p.add_run()
        r1.text = "—  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r1.font.name = font
        if isinstance(item, tuple):
            lead, rest = item
            r2 = p.add_run()
            r2.text = lead
            r2.font.size = Pt(size)
            r2.font.bold = True
            r2.font.color.rgb = color
            r2.font.name = font
            r3 = p.add_run()
            r3.text = rest
            r3.font.size = Pt(size)
            r3.font.bold = False
            r3.font.color.rgb = color
            r3.font.name = font
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = font
    return tb

def footer(s, idx, total, label):
    textbox(s, Inches(0.5), SH - Inches(0.45), Inches(6), Inches(0.35),
            f"Bina  —  بناء", size=10, color=MUTED, bold=True, font=FONT)
    textbox(s, Inches(4.6), SH - Inches(0.45), Inches(6), Inches(0.35),
            label, size=10, color=MUTED, align=PP_ALIGN.CENTER, font=FONT)
    textbox(s, SW - Inches(1.3), SH - Inches(0.45), Inches(0.8), Inches(0.35),
            f"{idx:02d} / {total:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)

def kicker_title(s, kicker, title, x=Inches(0.7), y=Inches(0.55), w=Inches(11.9),
                  title_size=32, kicker_color=ORANGE, title_color=NAVY):
    textbox(s, x, y, w, Inches(0.35), kicker.upper(), size=13, color=kicker_color, bold=True, font=FONT)
    textbox(s, x, y + Inches(0.36), w, Inches(0.9), title, size=title_size, color=title_color, bold=True, font=FONT)

TOTAL_SLIDES = 16

# ============================================================ SLIDE 1 — TITLE
s = add_slide(NAVY_DK)
rect(s, Inches(0), Inches(0), Inches(0.18), SH, ORANGE)
textbox(s, Inches(0.9), Inches(2.15), Inches(10), Inches(0.5), "PITCH DECK  ·  V0.1", size=14, color=ORANGE, bold=True, font=FONT)
tb = s.shapes.add_textbox(Inches(0.85), Inches(2.55), Inches(11.5), Inches(1.7))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Bina "; r.font.size = Pt(64); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
r2 = p.add_run(); r2.text = "بناء"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = ORANGE; r2.font.name = FONT_AR
textbox(s, Inches(0.9), Inches(4.05), Inches(10.8), Inches(0.9),
        "Trouvez les marchés. Formez votre équipe. Gagnez ensemble.",
        size=22, color=RGBColor(0xE8,0xE9,0xEC), italic=True, font=FONT)
textbox(s, Inches(0.9), Inches(4.65), Inches(10.8), Inches(0.5),
        "La première plateforme construite pour les PME du BTP marocain — veille des marchés publics, formation de groupements, et conformité réglementaire.",
        size=14, color=RGBColor(0xB9,0xC2,0xCF), font=FONT)
line = rect(s, Inches(0.9), Inches(5.55), Inches(3.2), Pt(2), ORANGE)
textbox(s, Inches(0.9), Inches(5.75), Inches(6), Inches(0.4), "bina.ma", size=13, color=WHITE, bold=True, font=FONT)
textbox(s, Inches(0.9), Inches(6.15), Inches(11), Inches(0.4),
        "Construit avec Claude Code · Données HCP, Médias24, Bati.ma, FNBTP", size=10.5, color=RGBColor(0x8B,0x96,0xA6), font=FONT)

# ============================================================ SLIDE 2 — THE PROBLEM (stats)
s = add_slide(CONCRETE)
kicker_title(s, "Le problème", "Le boom de la construction marocaine contourne les PME")
stats = [
    ("550 Md MAD", "d'investissements construction cumulés 2024–2030"),
    ("380 Md MAD", "investissement public — Loi de Finances 2026, un record historique"),
    ("1,24 M", "travailleurs BTP — 6% du PIB, +4,1% de croissance en 2026"),
    ("2,8 Md €", "impact économique projeté de la Coupe du Monde 2030 + 335 000 emplois"),
]
card_w, card_h, gap = Inches(2.85), Inches(2.0), Inches(0.22)
start_x = Inches(0.7)
y0 = Inches(1.75)
for i, (num, label) in enumerate(stats):
    x = start_x + i * (card_w + gap)
    c = rrect(s, x, y0, card_w, card_h, WHITE, radius=0.06, line=True, line_color=BORDER, line_w=Pt(0.75))
    rect(s, x, y0, card_w, Inches(0.08), ORANGE)
    textbox(s, x + Inches(0.2), y0 + Inches(0.35), card_w - Inches(0.4), Inches(0.7), num,
            size=30, color=NAVY, bold=True, font=FONT)
    textbox(s, x + Inches(0.2), y0 + Inches(1.15), card_w - Inches(0.4), Inches(0.75), label,
            size=12.5, color=MUTED, font=FONT, line_spacing=1.15)
textbox(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(0.4),
        "Mais 6 grands groupes captent l'essentiel des grands marchés. Les PME font face à 3 blocages :",
        size=15, color=INK, bold=True, font=FONT)
bullets(s, Inches(0.7), Inches(4.7), Inches(11.6), Inches(2.2), [
    ("marchespublics.gov.ma est inexploitable — ", "200+ AO/semaine, aucun filtrage intelligent ; 10–15h/semaine de veille manuelle."),
    ("Les gros marchés exigent des groupements — ", "qui ne se forment pas d'eux-mêmes ; aucun mécanisme pour trouver des partenaires complémentaires."),
    ("Les dossiers de conformité bloquent les PME — ", "attestation fiscale, quitus CNSS, assurance décennale... des jours à rassembler à chaque candidature."),
], size=14.5)
footer(s, 2, TOTAL_SLIDES, "Le problème")

# ============================================================ SLIDE 3 — POSITIONING
s = add_slide(NAVY)
rect(s, Inches(0), Inches(0), SW, Inches(0.12), ORANGE)
textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.4), "POSITIONNEMENT", size=13, color=ORANGE, bold=True, font=FONT)
tb = s.shapes.add_textbox(Inches(0.85), Inches(2.75), Inches(11.6), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.line_spacing = 1.25
r = p.add_run(); r.text = "marchespublics.gov.ma vous dit ce qui existe.\n"
r.font.size = Pt(30); r.font.color.rgb = RGBColor(0xC7,0xCE,0xD8); r.font.name = FONT; r.font.italic=True
p2 = tf.add_paragraph()
p2.line_spacing = 1.2
r2 = p2.add_run(); r2.text = "Bina vous dit ce que vous pouvez gagner —"
r2.font.size = Pt(34); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = FONT
p3 = tf.add_paragraph()
p3.line_spacing = 1.2
r3 = p3.add_run(); r3.text = "et vous connecte aux partenaires pour le gagner."
r3.font.size = Pt(34); r3.font.bold = True; r3.font.color.rgb = ORANGE; r3.font.name = FONT
footer(s, 3, TOTAL_SLIDES, "Positionnement")

# ============================================================ SLIDE 4 — SOLUTION OVERVIEW
s = add_slide(CONCRETE)
kicker_title(s, "La solution", "Trois modules, un seul flux de travail")
mods = [
    ("A", "Tender Radar", "Veille Marchés", "Agrège et filtre marchespublics.gov.ma : alertes intelligentes par métier, région, budget. Remplace 10–15h/semaine de veille manuelle."),
    ("B", "Groupement Builder", "Formation d'équipe", "Le différenciateur — rien de tel n'existe au Maroc. Trouvez des partenaires complémentaires et formez un groupement conforme à la loi."),
    ("C", "Compliance Kit", "Dossier Réglementaire", "Coffre-fort documentaire, suivi des expirations, générateur de dossier automatique. Des jours réduits à quelques minutes."),
]
card_w, card_h, gap = Inches(3.75), Inches(4.3), Inches(0.25)
x0 = Inches(0.7); y0 = Inches(1.85)
for i, (letter, name, sub, desc) in enumerate(mods):
    x = x0 + i * (card_w + gap)
    rrect(s, x, y0, card_w, card_h, WHITE, radius=0.045, line=True, line_color=BORDER, line_w=Pt(0.75))
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y0 + Inches(0.3), Inches(0.7), Inches(0.7))
    badge.fill.solid(); badge.fill.fore_color.rgb = NAVY; badge.line.fill.background(); badge.shadow.inherit=False
    tf = badge.text_frame; tf.word_wrap=False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = letter; r.font.size=Pt(26); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
    textbox(s, x + Inches(0.3), y0 + Inches(1.2), card_w - Inches(0.6), Inches(0.5), name, size=19, color=NAVY, bold=True, font=FONT)
    textbox(s, x + Inches(0.3), y0 + Inches(1.68), card_w - Inches(0.6), Inches(0.4), sub, size=12.5, color=ORANGE, bold=True, font=FONT)
    textbox(s, x + Inches(0.3), y0 + Inches(2.2), card_w - Inches(0.6), Inches(2.0), desc, size=13, color=MUTED, font=FONT, line_spacing=1.25)
footer(s, 4, TOTAL_SLIDES, "La solution")

# ============================================================ SLIDE 5 — MODULE A DETAIL
def module_detail_slide(idx, letter, title, subtitle, features, color=NAVY):
    s = add_slide(CONCRETE)
    rect(s, Inches(0), Inches(0), Inches(3.9), SH, color)
    textbox(s, Inches(0.6), Inches(0.6), Inches(1.2), Inches(1.0), letter, size=54, color=ORANGE, bold=True, font=FONT)
    textbox(s, Inches(0.6), Inches(1.75), Inches(3.0), Inches(1.4), title, size=27, color=WHITE, bold=True, font=FONT, line_spacing=1.1)
    textbox(s, Inches(0.6), Inches(3.05), Inches(3.0), Inches(0.6), subtitle, size=14, color=ORANGE, bold=True, font=FONT)
    textbox(s, Inches(0.6), SH - Inches(1.1), Inches(3.0), Inches(0.6), "Module " + letter, size=11, color=RGBColor(0xB9,0xC2,0xCF), font=FONT)
    bullets(s, Inches(4.35), Inches(0.85), Inches(8.3), Inches(6.0), features, size=15, gap=16, font=FONT)
    footer(s, idx, TOTAL_SLIDES, title)
    return s

module_detail_slide(5, "A", "Tender Radar", "Veille Marchés", [
    ("Flux agrégé — ", "scraping + parsing de marchespublics.gov.ma → base structurée : échéance, budget, région, maître d'ouvrage, lots, métiers requis."),
    ("Filtrage intelligent — ", "par métier (plomberie, électricité, génie civil...), région, budget (< 1M / 1–5M / 5–20M / > 20M MAD), type, maître d'ouvrage."),
    ("Recherches sauvegardées + alertes — ", "\"Nouvelle AO : plomberie, Casablanca, 2–8M MAD\" en email/in-app dès publication."),
    ("Compte à rebours d'échéance — ", "tableau de bord personnel, rouge < 7 j, orange < 14 j."),
    ("Détection de découpage en lots — ", "signale les AO scindées en lots qu'une PME peut remporter seule."),
])

# ============================================================ SLIDE 6 — MODULE B DETAIL
module_detail_slide(6, "B", "Groupement Builder", "Formation d'équipe", [
    ("Annonces de groupement — ", "\"Je cherche un partenaire électricité pour l'AO [X], ma part ~3M MAD, délai 6 mois — qui est intéressé ?\""),
    ("Annuaire des métiers — ", "PME vérifiées par spécialité, région et capacité (références)."),
    ("Suggestions de correspondance — ", "\"Cette AO exige 5 métiers. Vous couvrez la plomberie. 3 entreprises couvrent les autres lots. Former un groupement ?\""),
    ("Espace de travail du groupement — ", "documents partagés, répartition des tâches, suivi de contribution, messagerie."),
    ("Statuts suivis — ", "en formation → constitué → candidature déposée → résultat, avec un mandataire unique conforme au Décret 2-12-349."),
])

# ============================================================ SLIDE 7 — MODULE C DETAIL
module_detail_slide(7, "C", "Compliance Kit", "Dossier Réglementaire", [
    ("Coffre-fort documentaire — ", "attestation fiscale, quitus CNSS, déclaration sur l'honneur, assurance décennale, références — stockage privé R2, URLs signées 15 min."),
    ("Suivi d'expiration — ", "alertes 15 jours avant échéance ; aucun document expiré découvert au dépôt d'une candidature."),
    ("Score de conformité — ", "complétude du profil affichée de 0 à 100 %, visible avant de s'engager sur un marché."),
    ("Générateur de dossier — ", "sélectionnez un marché → liste des documents requis générée automatiquement, pré-remplie depuis le coffre-fort."),
    ("Aucune fausse garantie — ", "Bina organise les documents mais ne certifie jamais la conformité d'un contractant — disclaimer clair sur chaque dossier."),
])

# ============================================================ SLIDE 8 — WHY NOW (WC2030 timeline-ish)
s = add_slide(CONCRETE)
kicker_title(s, "Pourquoi maintenant", "La fenêtre Coupe du Monde 2030")
points = [
    ("96 Md MAD", "Ferroviaire seul — TGV + 150 trains + 40 nouvelles gares d'ici 2030"),
    ("3,8 Md MAD", "Prêt BAD pour les infrastructures Mondial, + 7 Md MAD de pipeline additionnel"),
    ("+3,8% AAGR", "Croissance prévisionnelle du secteur BTP 2026–2029"),
    ("335 000", "Emplois projetés par l'impact économique de la Coupe du Monde 2030"),
]
card_w, card_h, gap = Inches(2.85), Inches(1.85), Inches(0.22)
x0 = Inches(0.7); y0 = Inches(1.85)
for i, (num, label) in enumerate(points):
    x = x0 + i * (card_w + gap)
    rrect(s, x, y0, card_w, card_h, NAVY, radius=0.06)
    textbox(s, x + Inches(0.22), y0 + Inches(0.28), card_w - Inches(0.44), Inches(0.65), num, size=25, color=ORANGE, bold=True, font=FONT)
    textbox(s, x + Inches(0.22), y0 + Inches(0.95), card_w - Inches(0.44), Inches(0.8), label, size=11.5, color=RGBColor(0xE3,0xE7,0xEC), font=FONT, line_spacing=1.15)
textbox(s, Inches(0.7), Inches(4.15), Inches(11.9), Inches(1.9),
        "Ce cycle d'investissement historique se produit une fois par génération. Les PME qui structurent aujourd'hui "
        "leur veille, leurs partenariats et leur conformité seront celles qui captent la vague — les autres regarderont "
        "les 6 grands groupes absorber les marchés comme toujours.",
        size=16, color=INK, italic=True, line_spacing=1.35, font=FONT)
footer(s, 8, TOTAL_SLIDES, "Pourquoi maintenant")

# ============================================================ SLIDE 9 — DIFFERENTIATION TABLE
s = add_slide(CONCRETE)
kicker_title(s, "Différenciation", "Ce que fait Bina que le statu quo ne fait pas")
rows = [
    ("", "Statu quo (portail + réseau perso)", "Bina"),
    ("Filtrage des AO", "Manuel, chronophage (10–15h/sem)", "Automatique, par métier/région/budget"),
    ("Alertes", "Aucune — vérification manuelle quotidienne", "Email + in-app en temps réel"),
    ("Formation de groupement", "Réseau personnel, ad hoc, lent", "Annuaire vérifié + suggestions automatiques"),
    ("Dossier de conformité", "Reconstruit à chaque candidature", "Coffre-fort + générateur automatique"),
    ("Suivi des échéances", "Agenda personnel, risque d'oubli", "Compte à rebours visuel sur chaque AO"),
]
tbl_x, tbl_y, tbl_w, tbl_h = Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.7)
gtbl = s.shapes.add_table(len(rows), 3, tbl_x, tbl_y, tbl_w, tbl_h).table
gtbl.columns[0].width = Inches(3.1)
gtbl.columns[1].width = Inches(4.6)
gtbl.columns[2].width = Inches(4.2)
for ci in range(3):
    cell = gtbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = rows[0][ci]; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE; r.font.name = FONT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
for ri in range(1, len(rows)):
    for ci in range(3):
        cell = gtbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else RGBColor(0xF0,0xEE,0xEA)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = rows[ri][ci]
        r.font.size = Pt(12.5)
        r.font.name = FONT
        r.font.bold = (ci == 0)
        r.font.color.rgb = NAVY if ci == 0 else (ORANGE if ci == 2 else MUTED)
        if ci == 2:
            r.font.bold = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
# remove default table style borders look by setting a clean style via XML banding off
gtbl.first_row = False
gtbl.horz_banding = False
footer(s, 9, TOTAL_SLIDES, "Différenciation")

# ============================================================ SLIDE 10 — ARCHITECTURE / TECH CREDIBILITY
s = add_slide(NAVY_DK)
kicker_title(s, "Fondations techniques", "Construit pour la fiabilité et la conformité", title_color=WHITE)
left_items = [
    ("Next.js 15 + TypeScript strict — ", "rendu SSR des pages publiques d'AO pour le SEO organique."),
    ("PostgreSQL 16 + Drizzle + RLS — ", "isolation stricte des rôles au niveau base de données."),
    ("Playwright scraper — ", "respecte marchespublics.gov.ma : 1 requête/3s, nightly only, fallback CSV manuel."),
    ("Cloudflare R2 privé — ", "documents de conformité, URLs signées 15 min, accès audité."),
]
right_items = [
    ("Deux rôles seulement — ", "contractant / admin. Pas de flux d'argent — Bina est un outil SaaS, jamais intermédiaire."),
    ("i18n FR/AR — ", "RTL obligatoire et traité à égalité, pas en traduction de second rang."),
    ("Couverture de tests ≥ 80% — ", "Vitest + Playwright, seuil appliqué en CI, jamais abaissé."),
    ("CI/CD GitHub Actions — ", "chaque push est vérifié, aucune fonctionnalité n'est \"terminée\" tant que le pipeline n'est pas vert."),
]
textbox(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.4), "Produit & Données", size=13, color=ORANGE, bold=True, font=FONT)
bullets(s, Inches(0.7), Inches(2.15), Inches(5.6), Inches(4.5), left_items, size=13.5, color=RGBColor(0xE3,0xE7,0xEC), gap=14, font=FONT)
textbox(s, Inches(6.85), Inches(1.7), Inches(5.6), Inches(0.4), "Sécurité & Qualité", size=13, color=ORANGE, bold=True, font=FONT)
bullets(s, Inches(6.85), Inches(2.15), Inches(5.6), Inches(4.5), right_items, size=13.5, color=RGBColor(0xE3,0xE7,0xEC), gap=14, font=FONT)
footer(s, 10, TOTAL_SLIDES, "Fondations techniques")

# ============================================================ SLIDE 11 — DESIGN IDENTITY
s = add_slide(CONCRETE)
kicker_title(s, "Identité de marque", "Ton pratique, ni startup ni portail administratif")
textbox(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.4), "Palette", size=13, color=ORANGE, bold=True, font=FONT)
swatches = [("Steel Blue", NAVY, "#1E3A5F"), ("Construction Orange", ORANGE, "#E07B39"), ("Concrete", CONCRETE, "#F4F2EF")]
sx = Inches(0.7)
for name, col, hexv in swatches:
    rrect(s, sx, Inches(2.35), Inches(1.7), Inches(1.7), col, radius=0.12, line=True, line_color=BORDER, line_w=Pt(0.75))
    textbox(s, sx, Inches(4.15), Inches(1.7), Inches(0.35), name, size=11, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT)
    textbox(s, sx, Inches(4.5), Inches(1.7), Inches(0.3), hexv, size=10, color=MUTED, align=PP_ALIGN.CENTER, font=FONT)
    sx += Inches(1.95)
textbox(s, Inches(0.7), Inches(5.15), Inches(5.7), Inches(0.4), "Typographie", size=13, color=ORANGE, bold=True, font=FONT)
textbox(s, Inches(0.7), Inches(5.6), Inches(5.7), Inches(0.5), "IBM Plex Sans", size=20, bold=True, color=NAVY, font=FONT)
textbox(s, Inches(0.7), Inches(6.15), Inches(5.7), Inches(0.5), "بناء — Noto Kufi Arabic", size=18, bold=True, color=NAVY, font=FONT_AR)
textbox(s, Inches(7.1), Inches(1.85), Inches(5.5), Inches(0.4), "Voix éditoriale", size=13, color=ORANGE, bold=True, font=FONT)
bullets(s, Inches(7.1), Inches(2.35), Inches(5.5), Inches(4.4), [
    "Pratique, professionnel, entre pairs — l'énergie d'un confrère du BTP, pas d'un discours startup.",
    "Ni le hype de la tech, ni la rigidité d'un portail gouvernemental.",
    "« Bina » — بناء, \"construction\" en darija : compris immédiatement par tout le secteur.",
    "Signal d'urgence constant : compte à rebours d'échéance visible partout, rouge < 7 j.",
], size=13.5, gap=16)
footer(s, 11, TOTAL_SLIDES, "Identité de marque")

# ============================================================ SLIDE 12 — BUSINESS MODEL
s = add_slide(CONCRETE)
kicker_title(s, "Modèle économique", "SaaS par abonnement — jamais intermédiaire financier")
bullets(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.5), [
    ("Abonnement SaaS mensuel — ", "facturé par compte contractant ; paliers par taille d'entreprise (à définir selon pilote)."),
    ("Aucune commission sur marché — ", "Bina ne prend jamais de marge sur un contrat gagné — outil, pas agent de marché."),
    ("Aucun flux d'argent entre PME — ", "les transactions de groupement restent hors plateforme, conformément au cadre légal."),
    ("Acquisition organique — ", "les pages publiques d'AO (SSR, sans compte) servent de moteur SEO — le premier contact est gratuit et utile."),
], size=14.5, gap=16)
rrect(s, Inches(7.1), Inches(1.9), Inches(5.55), Inches(4.5), NAVY, radius=0.05)
textbox(s, Inches(7.4), Inches(2.15), Inches(5.0), Inches(0.4), "CE QUE BINA NE FAIT JAMAIS", size=12, color=ORANGE, bold=True, font=FONT)
bullets(s, Inches(7.4), Inches(2.65), Inches(5.0), Inches(3.5), [
    "Ne certifie jamais qu'un contractant est conforme.",
    "Ne fait jamais transiter d'argent de marché.",
    "Ne prend jamais de commission sur une candidature gagnée.",
    "Ne remplace jamais le maître d'ouvrage dans l'évaluation.",
], size=13, color=RGBColor(0xE3,0xE7,0xEC), bullet_color=ORANGE, gap=14)
footer(s, 12, TOTAL_SLIDES, "Modèle économique")

# ============================================================ SLIDE 13 — ROADMAP / STATUS
s = add_slide(CONCRETE)
kicker_title(s, "Feuille de route", "v0.1 livré — sprint par sprint")
sprints = [
    ("S0", "Scaffold, Auth, RBAC, RLS, Docker"),
    ("S1", "Modèle de données, profils, références"),
    ("S2", "Scraper + veille des marchés (SSR)"),
    ("S3", "Alertes + suivi des échéances"),
    ("S4", "Groupements — création à espace de travail"),
    ("S5", "Coffre-fort conformité + générateur dossier"),
    ("S6", "Notifications + i18n FR/AR + RTL"),
    ("S7", "Admin, durcissement sécurité, déploiement"),
]
n = len(sprints)
box_w = Inches(1.36)
gap = Inches(0.12)
total_w = box_w * n + gap * (n - 1)
x0 = (SW - total_w) / 2
y0 = Inches(2.5)
liney = y0 + Inches(0.35)
rect(s, x0, liney, total_w, Pt(2.5), GREEN)
for i, (tag, desc) in enumerate(sprints):
    x = x0 + i * (box_w + gap)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + box_w/2 - Inches(0.14), liney - Inches(0.12), Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.color.rgb=WHITE; dot.line.width=Pt(1.5); dot.shadow.inherit=False
    textbox(s, x, y0 - Inches(0.5), box_w, Inches(0.35), tag, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font=FONT)
    textbox(s, x - Inches(0.15), liney + Inches(0.35), box_w + Inches(0.3), Inches(1.3), desc, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, font=FONT, line_spacing=1.1)
textbox(s, x0, liney + Inches(1.85), total_w, Inches(0.4), "✓ TOUS LIVRÉS — v0.1 SHIPPED", size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER, font=FONT)

textbox(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.4), "PROCHAIN — v0.2 backlog", size=13, color=ORANGE, bold=True, font=FONT)
bullets(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.4), [
    "Application mobile (React Native)  ·  Assistant de rédaction d'AO par IA  ·  Marketplace sous-traitants & matériel",
    "Intégration financement Crédit Agricole du Maroc  ·  Vérification temps réel DGI/CNSS  ·  v0.3 : marchés internationaux (BERD, BAD, IDA)",
], size=12.5, gap=8)
footer(s, 13, TOTAL_SLIDES, "Feuille de route")

# ============================================================ SLIDE 14 — DEFINITION OF DONE / QUALITY BAR
s = add_slide(NAVY_DK)
kicker_title(s, "Statut v0.1", "20/20 critères de \"Definition of Done\" livrés", title_color=WHITE)
dod_left = [
    "Auth contractant/admin + vérification email",
    "Profil contractant : spécialités, régions, FNBTP",
    "Scraper : base peuplée (40+ AO réelles/mock)",
    "Veille publique filtrable, SSR",
    "Détail AO structuré + échéance + lots",
    "Recherches sauvegardées + alertes email",
    "Suivi de candidature (watching/bidding/soumis)",
    "Création de groupement par AO",
    "Annuaire des groupements par métier",
    "Rejoindre / quitter un groupement",
]
dod_right = [
    "Espace de travail groupement partagé",
    "Coffre-fort de documents (R2 privé)",
    "Suivi d'expiration + alertes 15 jours",
    "Score de conformité 0–100%",
    "Générateur automatique de dossier",
    "Références de projets (profil public)",
    "Notifications in-app",
    "Emails transactionnels (Resend)",
    "Dashboard admin : scraper, KPIs, groupements",
    "FR + AR + RTL complet · build/tests/lint verts",
]
def dod_col(x, items):
    tb = s.shapes.add_textbox(x, Inches(1.7), Inches(5.7), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r1 = p.add_run(); r1.text = "✓  "; r1.font.size=Pt(13); r1.font.bold=True; r1.font.color.rgb=GREEN; r1.font.name=FONT
        r2 = p.add_run(); r2.text = it; r2.font.size=Pt(12.5); r2.font.color.rgb=RGBColor(0xE3,0xE7,0xEC); r2.font.name=FONT
dod_col(Inches(0.7), dod_left)
dod_col(Inches(6.6), dod_right)
footer(s, 14, TOTAL_SLIDES, "Statut v0.1")

# ============================================================ SLIDE 15 — THE ASK
s = add_slide(CONCRETE)
rect(s, Inches(0), Inches(0), SW, Inches(0.12), ORANGE)
kicker_title(s, "L'invitation", "Rejoignez le lancement de Bina")
bullets(s, Inches(0.7), Inches(2.0), Inches(11.7), Inches(3.5), [
    ("PME du BTP — ", "rejoignez le programme pilote. Créez votre profil, activez vos alertes, testez le générateur de dossier sur votre prochaine candidature."),
    ("Associations professionnelles & FNBTP — ", "devenez partenaire de distribution auprès de vos adhérents — Bina structure ce que vos membres font déjà manuellement."),
    ("Investisseurs & partenaires — ", "le v0.1 est livré et testé ; la fenêtre Coupe du Monde 2030 est ouverte maintenant, pas dans deux ans."),
], size=16, gap=22)
rrect(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.1), NAVY, radius=0.08)
textbox(s, Inches(1.0), Inches(5.95), Inches(11.2), Inches(0.6),
        "bina.ma  ·  hassan.plomberie@demo.bina.ma — accès démo disponible sur demande",
        size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=FONT)
footer(s, 15, TOTAL_SLIDES, "L'invitation")

# ============================================================ SLIDE 16 — CLOSING
s = add_slide(NAVY_DK)
rect(s, Inches(0), Inches(0), Inches(0.18), SH, ORANGE)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(10.5), Inches(1.4))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Bina "; r.font.size=Pt(50); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
r2 = p.add_run(); r2.text = "بناء"; r2.font.size=Pt(42); r2.font.bold=True; r2.font.color.rgb=ORANGE; r2.font.name=FONT_AR
textbox(s, Inches(0.9), Inches(3.95), Inches(10.5), Inches(0.6),
        "Trouvez les marchés. Formez votre équipe. Gagnez ensemble.", size=18, color=RGBColor(0xC7,0xCE,0xD8), italic=True, font=FONT)
textbox(s, Inches(0.9), Inches(4.75), Inches(10.5), Inches(0.4), "bina.ma", size=15, color=ORANGE, bold=True, font=FONT)
textbox(s, Inches(0.9), Inches(6.7), Inches(10.5), Inches(0.4),
        "Construit avec Claude Code · Sprint 0–7 · v0.1 shipped", size=10.5, color=RGBColor(0x8B,0x96,0xA6), font=FONT)

out = "bina-pitch-deck.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides" if False else len(prs.slides._sldIdLst))
